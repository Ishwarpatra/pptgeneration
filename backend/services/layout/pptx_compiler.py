"""
PPTX Compiler - Generates PowerPoint files from structured content.
Uses python-pptx to create professional presentations.
"""
import os
import io
import logging
import httpx
from pathlib import Path
from typing import List, Optional, Union
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from services.nlp.models import PresentationOutline, SlideContent, LayoutType, VisualType, VisualSpec
from services.style.style_gene import StyleGene
from services.style.presets import get_preset, MODERN_MINIMAL
from services.config import (
    AspectRatio,
    get_slide_dimensions,
    DEFAULT_ASPECT_RATIO,
    SlideDimensions
)
from .paginator import Paginator, SlideSpec

# Configure logging
logger = logging.getLogger(__name__)


class PPTXCompiler:
    """
    Compiles presentation outlines into PowerPoint files.
    Applies Style Genes to formatting.
    """

    def __init__(
        self,
        style: Optional[StyleGene] = None,
        output_dir: str = "presentations",
        aspect_ratio: Optional[AspectRatio] = None,
        images_dir: str = "generated_images"
    ):
        """
        Initialize the compiler.
        
        Args:
            style: StyleGene to apply (defaults to Modern Minimal)
            output_dir: Directory to save generated files
            aspect_ratio: Slide aspect ratio (defaults to 16:9 widescreen)
            images_dir: Directory containing generated images
        """
        self.style = style or MODERN_MINIMAL
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.paginator = Paginator()
        
        # Configurable slide dimensions
        self.slide_dimensions = get_slide_dimensions(aspect_ratio)
        
        # Directory for generated images
        self.images_dir = Path(images_dir)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip("#")
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _apply_style_to_text(self, text_frame, is_title: bool = False):
        """Apply style gene formatting to a text frame."""
        colors = self.style.palette.to_hex_dict()
        typo = self.style.typography
        sizes = typo.get_heading_sizes()
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Font family
                run.font.name = typo.heading_font if is_title else typo.body_font
                
                # Font size
                if is_title:
                    run.font.size = Pt(sizes["h1"])
                else:
                    run.font.size = Pt(sizes["body"])
                
                # Color
                run.font.color.rgb = self._hex_to_rgb(colors["text_primary"])

    def _get_layout_index(self, layout_type: LayoutType, prs: Presentation) -> int:
        """Map LayoutType to PowerPoint layout index."""
        # Standard PowerPoint layout indices:
        # 0 = Title Slide
        # 1 = Title and Content
        # 2 = Section Header
        # 3 = Two Content
        # 4 = Comparison
        # 5 = Title Only
        # 6 = Blank
        
        layout_map = {
            LayoutType.TITLE_ONLY: 5,
            LayoutType.TITLE_CONTENT: 1,
            LayoutType.TWO_COLUMN: 3,
            LayoutType.SECTION_HEADER: 2,
            LayoutType.VISUAL_HEAVY: 1,
            LayoutType.COMPARISON: 4,
            LayoutType.BLANK: 6,
        }
        
        index = layout_map.get(layout_type, 1)
        # Ensure the layout exists
        if index >= len(prs.slide_layouts):
            return 1  # Fallback to Title and Content
        return index

    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: Optional[str] = None
    ):
        """Add the title slide."""
        layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_style_to_text(slide.shapes.title.text_frame, is_title=True)
        
        # Set subtitle
        if subtitle and len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
            self._apply_style_to_text(subtitle_placeholder.text_frame)

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_content: SlideContent,
        is_continuation: bool = False
    ):
        """Add a content slide."""
        layout_idx = self._get_layout_index(slide_content.layout_type, prs)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        # Set title
        title_text = slide_content.title
        if is_continuation:
            title_text += " (Continued)"
        
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            self._apply_style_to_text(slide.shapes.title.text_frame, is_title=True)
        
        # Find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break
        
        if content_placeholder and slide_content.body_text:
            tf = content_placeholder.text_frame
            tf.clear()
            
            for i, point in enumerate(slide_content.body_text):
                if i == 0:
                    tf.paragraphs[0].text = point
                    tf.paragraphs[0].level = 0
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
            
            self._apply_style_to_text(tf)
        
        # Add speaker notes
        if slide_content.speaker_notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_content.speaker_notes

    def _load_image_from_path(self, image_path: str) -> Optional[io.BytesIO]:
        """
        Load an image from a local file path.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            BytesIO stream of the image, or None if loading fails
        """
        try:
            path = Path(image_path)
            
            # Check if it's a relative path within the images directory
            if not path.is_absolute():
                path = self.images_dir / path
            
            if path.exists() and path.is_file():
                with open(path, "rb") as f:
                    return io.BytesIO(f.read())
            else:
                logger.warning(f"Image file not found: {path}")
                return None
        except (IOError, OSError) as e:
            logger.error(f"Failed to load image from path {image_path}: {e}")
            return None

    def _download_image_from_url(self, url: str) -> Optional[io.BytesIO]:
        """
        Download an image from a URL.
        
        Args:
            url: URL of the image
            
        Returns:
            BytesIO stream of the image, or None if download fails
        """
        try:
            with httpx.Client(timeout=30.0) as client:
                response = client.get(url)
                response.raise_for_status()
                return io.BytesIO(response.content)
        except httpx.HTTPError as e:
            logger.error(f"Failed to download image from URL {url}: {e}")
            return None
        except Exception as e:
            logger.error(f"Unexpected error downloading image from {url}: {e}")
            return None

    def _get_image_stream(self, visual_spec: VisualSpec) -> Optional[io.BytesIO]:
        """
        Get image stream from visual spec, trying path first, then URL.
        
        Args:
            visual_spec: The visual specification containing path or URL
            
        Returns:
            BytesIO stream of the image, or None if not available
        """
        # First, check if there's a local image path in the prompt or metadata
        # The visual engine stores the path in the image_path field
        image_path = getattr(visual_spec, 'image_path', None)
        
        if image_path:
            stream = self._load_image_from_path(image_path)
            if stream:
                return stream
        
        # Check if the prompt contains a file path (for already-generated images)
        prompt = visual_spec.prompt
        if prompt and (prompt.endswith('.png') or prompt.endswith('.jpg') or prompt.endswith('.webp')):
            stream = self._load_image_from_path(prompt)
            if stream:
                return stream
        
        # Check if prompt is a URL
        if prompt and prompt.startswith(('http://', 'https://')):
            stream = self._download_image_from_url(prompt)
            if stream:
                return stream
        
        return None

    def _add_visual_element(
        self,
        slide,
        visual_spec: VisualSpec
    ) -> bool:
        """
        Add a visual element (image) to the slide.
        
        Attempts to embed an actual image. Falls back to placeholder if
        no image is available.
        
        Args:
            slide: The slide to add the visual to
            visual_spec: The visual specification
            
        Returns:
            True if an actual image was embedded, False if a placeholder was used
        """
        if not visual_spec or visual_spec.visual_type == VisualType.NONE:
            return False
        
        colors = self.style.palette.to_hex_dict()
        
        # Visual positioning (right side of slide)
        left = Inches(self.slide_dimensions.width * 0.55)
        top = Inches(1.5)
        width = Inches(self.slide_dimensions.width * 0.4)
        height = Inches(self.slide_dimensions.height * 0.6)
        
        # Attempt to get actual image
        image_stream = self._get_image_stream(visual_spec)
        
        if image_stream:
            try:
                # Add the actual image to the slide
                slide.shapes.add_picture(
                    image_stream,
                    left, top, width, height
                )
                logger.info(f"Successfully embedded image for visual: {visual_spec.visual_type.value}")
                return True
            except Exception as e:
                logger.error(f"Failed to embed image: {e}")
                # Fall through to placeholder
        
        # Fallback: Create a styled placeholder
        logger.warning(
            f"No image available for visual '{visual_spec.visual_type.value}', "
            f"creating placeholder. Prompt: {visual_spec.prompt[:50]}..."
        )
        self._add_visual_placeholder(slide, visual_spec, left, top, width, height)
        return False

    def _add_visual_placeholder(
        self,
        slide,
        visual_spec: VisualSpec,
        left: float,
        top: float,
        width: float,
        height: float
    ):
        """
        Add a styled placeholder rectangle when no actual image is available.
        
        This is a fallback that clearly indicates an image should be here.
        """
        colors = self.style.palette.to_hex_dict()
        
        # Create a placeholder rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        
        # Style the placeholder with a distinctive look
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(colors["surface"])
        shape.line.color.rgb = self._hex_to_rgb(colors["secondary"])
        shape.line.width = Pt(2)
        
        # Add informative label
        tf = shape.text_frame
        tf.word_wrap = True
        
        # Main indicator
        p = tf.paragraphs[0]
        p.text = f"📷 [{visual_spec.visual_type.value.upper()}]"
        p.alignment = PP_ALIGN.CENTER
        
        # Add description of what should be here
        p2 = tf.add_paragraph()
        prompt_preview = visual_spec.prompt[:80] + "..." if len(visual_spec.prompt) > 80 else visual_spec.prompt
        p2.text = f"\n{prompt_preview}"
        p2.alignment = PP_ALIGN.CENTER
        
        # Add note about generation
        p3 = tf.add_paragraph()
        p3.text = "\n[Image pending generation]"
        p3.alignment = PP_ALIGN.CENTER

    def compile(
        self,
        outline: PresentationOutline,
        filename: Optional[str] = None,
        aspect_ratio: Optional[AspectRatio] = None
    ) -> str:
        """
        Compile a presentation outline to a PPTX file.
        
        Args:
            outline: The structured presentation outline
            filename: Optional output filename (auto-generated if not provided)
            aspect_ratio: Override the default aspect ratio for this presentation
            
        Returns:
            Path to the generated PPTX file
        """
        prs = Presentation()
        
        # Use provided aspect ratio or fall back to instance setting
        dimensions = get_slide_dimensions(aspect_ratio) if aspect_ratio else self.slide_dimensions
        
        # Set slide dimensions from configuration
        prs.slide_width = Inches(dimensions.width)
        prs.slide_height = Inches(dimensions.height)
        
        logger.info(
            f"Creating presentation with {dimensions.aspect_ratio.value} aspect ratio "
            f"({dimensions.width}\" x {dimensions.height}\")"
        )
        
        # Add title slide
        self._add_title_slide(prs, outline.title, outline.subtitle)
        
        # Paginate and add content slides
        paginated_slides = self.paginator.paginate_presentation(outline.slides)
        
        images_embedded = 0
        images_placeholder = 0
        
        for i, slide_spec in enumerate(outline.slides):
            self._add_content_slide(prs, slide_spec)
            
            # Add visual element if specified
            if slide_spec.visual_spec and slide_spec.visual_spec.visual_type != VisualType.NONE:
                slide = prs.slides[-1]  # Get the just-added slide
                if self._add_visual_element(slide, slide_spec.visual_spec):
                    images_embedded += 1
                else:
                    images_placeholder += 1
        
        logger.info(
            f"Presentation compiled: {images_embedded} images embedded, "
            f"{images_placeholder} placeholders created"
        )
        
        # Generate filename
        if not filename:
            safe_title = "".join(c if c.isalnum() else "_" for c in outline.title)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{safe_title[:50]}_{timestamp}.pptx"
        
        # Save the file
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        
        return str(output_path)

    def compile_with_style(
        self,
        outline: PresentationOutline,
        style_id: str,
        filename: Optional[str] = None,
        aspect_ratio: Optional[AspectRatio] = None
    ) -> str:
        """
        Compile using a specific style preset.
        
        Args:
            outline: The structured presentation outline
            style_id: ID of the style preset to use
            filename: Optional output filename
            aspect_ratio: Optional aspect ratio override
            
        Returns:
            Path to generated file
        """
        self.style = get_preset(style_id)
        return self.compile(outline, filename, aspect_ratio)


# Convenience function
def generate_pptx(
    outline: PresentationOutline,
    style_id: str = "modern_minimal",
    aspect_ratio: Optional[AspectRatio] = None
) -> str:
    """
    Quick function to generate a PPTX from an outline.
    
    Args:
        outline: Structured presentation outline
        style_id: Style preset to apply
        aspect_ratio: Optional aspect ratio (defaults to 16:9)
        
    Returns:
        Path to generated file
    """
    compiler = PPTXCompiler(aspect_ratio=aspect_ratio)
    return compiler.compile_with_style(outline, style_id, aspect_ratio=aspect_ratio)
