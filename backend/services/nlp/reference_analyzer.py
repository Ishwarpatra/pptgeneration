"""
Reference PPT Analyzer - Extracts Style Genes from uploaded presentations.
Combines XML parsing with visual analysis for comprehensive style extraction.
"""
import io
import os
import tempfile
import logging
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field
from collections import Counter
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Configure logging for this module
logger = logging.getLogger(__name__)

from services.style.style_gene import (
    StyleGene,
    TypographyParams,
    ColorPalette,
    LayoutPhysics,
    LABColor,
    FontWeight,
)


@dataclass
class ExtractedColor:
    """A color extracted from the presentation with usage context."""
    hex_value: str
    rgb: Tuple[int, int, int]
    count: int = 1
    context: str = ""  # "background", "text", "accent", etc.


@dataclass
class ExtractedFont:
    """A font extracted from the presentation."""
    name: str
    count: int = 1
    is_heading: bool = False
    avg_size_pt: float = 0.0


@dataclass
class LayoutMetrics:
    """Metrics about slide layout density and structure."""
    avg_shapes_per_slide: float = 0.0
    avg_text_boxes: float = 0.0
    avg_images: float = 0.0
    content_density: float = 0.0  # 0-1 scale
    avg_margin_ratio: float = 0.0


@dataclass
class ExtractionResult:
    """Complete extraction result from a reference presentation."""
    colors: List[ExtractedColor] = field(default_factory=list)
    fonts: List[ExtractedFont] = field(default_factory=list)
    layout_metrics: LayoutMetrics = field(default_factory=LayoutMetrics)
    slide_count: int = 0
    has_master_theme: bool = False
    theme_name: Optional[str] = None


class ReferenceAnalyzer:
    """
    Analyzes reference PowerPoint files to extract Style Genes.
    Uses both XML parsing and visual analysis techniques.
    """

    def __init__(self):
        self.color_counter = Counter()
        self.font_counter = Counter()
        self.heading_fonts = Counter()
        self.body_fonts = Counter()

    def _rgb_to_hex(self, r: int, g: int, b: int) -> str:
        """Convert RGB values to hex string."""
        return f"#{r:02x}{g:02x}{b:02x}"

    def _extract_color_from_rgb(self, rgb_color) -> Optional[ExtractedColor]:
        """Extract color from python-pptx RGBColor object."""
        if rgb_color is None:
            return None
        try:
            r = (rgb_color >> 16) & 0xFF
            g = (rgb_color >> 8) & 0xFF
            b = rgb_color & 0xFF
            hex_val = self._rgb_to_hex(r, g, b)
            return ExtractedColor(
                hex_value=hex_val,
                rgb=(r, g, b),
            )
        except (TypeError, AttributeError):
            return None

    def _analyze_shape_colors(self, shape, context: str = "") -> List[ExtractedColor]:
        """Extract colors from a single shape."""
        colors = []
        
        # Extract fill color
        try:
            if hasattr(shape, 'fill') and shape.fill is not None:
                fill = shape.fill
                if hasattr(fill, 'fore_color') and fill.fore_color is not None:
                    fc = fill.fore_color
                    if hasattr(fc, 'rgb') and fc.rgb is not None:
                        rgb = fc.rgb
                        hex_val = self._rgb_to_hex(rgb[0], rgb[1], rgb[2])
                        colors.append(ExtractedColor(
                            hex_value=hex_val,
                            rgb=(rgb[0], rgb[1], rgb[2]),
                            context="fill"
                        ))
        except (AttributeError, TypeError) as e:
            logger.debug(f"Could not extract fill color from shape: {e}")
        except (IndexError, ValueError) as e:
            logger.warning(f"Invalid fill color data in shape: {e}")
        
        # Extract text colors
        try:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, 'color') and run.font.color is not None:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                                rgb = run.font.color.rgb
                                hex_val = self._rgb_to_hex(rgb[0], rgb[1], rgb[2])
                                colors.append(ExtractedColor(
                                    hex_value=hex_val,
                                    rgb=(rgb[0], rgb[1], rgb[2]),
                                    context="text"
                                ))
        except (AttributeError, TypeError) as e:
            logger.debug(f"Could not extract text color from shape: {e}")
        except (IndexError, ValueError) as e:
            logger.warning(f"Invalid text color data in shape: {e}")
        
        return colors

    def _analyze_shape_fonts(self, shape, is_title: bool = False) -> List[ExtractedFont]:
        """Extract font information from a shape."""
        fonts = []
        
        try:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, 'name') and run.font.name:
                            size = 0.0
                            try:
                                if hasattr(run.font, 'size') and run.font.size:
                                    size = run.font.size.pt if run.font.size else 0.0
                            except (AttributeError, TypeError) as e:
                                logger.debug(f"Could not read font size: {e}")
                                size = 0.0  # Default to 0 when size cannot be read
                            
                            fonts.append(ExtractedFont(
                                name=run.font.name,
                                is_heading=is_title or size > 24,
                                avg_size_pt=size
                            ))
        except (AttributeError, TypeError) as e:
            logger.debug(f"Could not extract fonts from shape: {e}")
        except Exception as e:
            # Catch-all for unexpected errors, but log them
            logger.error(f"Unexpected error extracting fonts from shape: {type(e).__name__}: {e}")
        
        return fonts

    def _calculate_layout_metrics(self, prs: Presentation) -> LayoutMetrics:
        """Calculate layout density and structure metrics."""
        total_shapes = 0
        total_text_boxes = 0
        total_images = 0
        slide_count = len(prs.slides)
        
        if slide_count == 0:
            return LayoutMetrics()
        
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_area = slide_width * slide_height if slide_width and slide_height else 1
        
        total_content_area = 0
        margin_samples = []
        
        for slide in prs.slides:
            shapes_on_slide = 0
            content_area_on_slide = 0
            min_left = float('inf')
            min_top = float('inf')
            
            for shape in slide.shapes:
                shapes_on_slide += 1
                total_shapes += 1
                
                # Count shape types
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    total_text_boxes += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    total_images += 1
                
                # Calculate content area
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    if shape.width and shape.height:
                        content_area_on_slide += shape.width * shape.height
                
                # Track margins
                if hasattr(shape, 'left') and shape.left:
                    min_left = min(min_left, shape.left)
                if hasattr(shape, 'top') and shape.top:
                    min_top = min(min_top, shape.top)
            
            total_content_area += content_area_on_slide
            
            # Estimate margin ratio
            if min_left != float('inf') and slide_width:
                margin_samples.append(min_left / slide_width)
        
        avg_density = (total_content_area / slide_count) / slide_area if slide_area else 0
        avg_margin = sum(margin_samples) / len(margin_samples) if margin_samples else 0.08
        
        return LayoutMetrics(
            avg_shapes_per_slide=total_shapes / slide_count,
            avg_text_boxes=total_text_boxes / slide_count,
            avg_images=total_images / slide_count,
            content_density=min(1.0, avg_density),
            avg_margin_ratio=avg_margin,
        )

    def analyze_presentation(self, file_path: str) -> ExtractionResult:
        """
        Analyze a PowerPoint file and extract all style information.
        
        Args:
            file_path: Path to the .pptx file
            
        Returns:
            ExtractionResult with all extracted style data
        """
        prs = Presentation(file_path)
        
        all_colors = []
        all_fonts = []
        
        # Analyze each slide
        for slide in prs.slides:
            # Check if this slide has a title shape
            title_shape = slide.shapes.title if hasattr(slide.shapes, 'title') else None
            
            for shape in slide.shapes:
                is_title = (shape == title_shape) if title_shape else False
                
                # Extract colors
                colors = self._analyze_shape_colors(shape)
                all_colors.extend(colors)
                
                # Extract fonts
                fonts = self._analyze_shape_fonts(shape, is_title)
                all_fonts.extend(fonts)
        
        # Aggregate colors by frequency
        color_counts = Counter()
        for c in all_colors:
            color_counts[c.hex_value] += 1
        
        aggregated_colors = [
            ExtractedColor(
                hex_value=hex_val,
                rgb=tuple(int(hex_val.lstrip('#')[i:i+2], 16) for i in (0, 2, 4)),
                count=count,
            )
            for hex_val, count in color_counts.most_common(20)
        ]
        
        # Aggregate fonts
        heading_fonts = Counter()
        body_fonts = Counter()
        font_sizes = {}
        
        for f in all_fonts:
            if f.is_heading:
                heading_fonts[f.name] += 1
            else:
                body_fonts[f.name] += 1
            
            if f.name not in font_sizes:
                font_sizes[f.name] = []
            if f.avg_size_pt > 0:
                font_sizes[f.name].append(f.avg_size_pt)
        
        aggregated_fonts = []
        all_font_names = set(heading_fonts.keys()) | set(body_fonts.keys())
        for name in all_font_names:
            is_heading = heading_fonts[name] > body_fonts[name]
            count = heading_fonts[name] + body_fonts[name]
            avg_size = sum(font_sizes.get(name, [12])) / len(font_sizes.get(name, [12]))
            aggregated_fonts.append(ExtractedFont(
                name=name,
                count=count,
                is_heading=is_heading,
                avg_size_pt=avg_size,
            ))
        
        # Sort by count
        aggregated_fonts.sort(key=lambda x: x.count, reverse=True)
        
        # Calculate layout metrics
        layout_metrics = self._calculate_layout_metrics(prs)
        
        return ExtractionResult(
            colors=aggregated_colors,
            fonts=aggregated_fonts[:10],  # Top 10 fonts
            layout_metrics=layout_metrics,
            slide_count=len(prs.slides),
            has_master_theme=len(prs.slide_masters) > 0,
        )

    def analyze_from_bytes(self, file_bytes: bytes) -> ExtractionResult:
        """
        Analyze a PowerPoint file from bytes (for file uploads).
        """
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        
        try:
            result = self.analyze_presentation(tmp_path)
        finally:
            os.unlink(tmp_path)
        
        return result

    def extract_style_gene(
        self,
        extraction_result: ExtractionResult,
        gene_name: str = "extracted_style"
    ) -> StyleGene:
        """
        Convert extraction results into a StyleGene.
        
        Args:
            extraction_result: The extraction result from analyze_presentation
            gene_name: Name for the generated StyleGene
            
        Returns:
            A StyleGene that approximates the reference presentation's style
        """
        # Determine primary colors
        colors = extraction_result.colors
        
        # Filter out very common colors (white, black, gray)
        filtered_colors = [
            c for c in colors 
            if c.hex_value.lower() not in ['#ffffff', '#000000', '#f0f0f0', '#333333']
        ]
        
        # Use filtered if available, otherwise use all
        color_source = filtered_colors if filtered_colors else colors
        
        # Build color palette
        def get_color_or_default(index: int, default: str) -> str:
            if index < len(color_source):
                return color_source[index].hex_value
            return default
        
        # Find background color (usually white or very light)
        background = "#ffffff"
        for c in colors:
            r, g, b = c.rgb
            if (r + g + b) / 3 > 240:  # Very light color
                background = c.hex_value
                break
        
        # Find text color (usually dark)
        text_color = "#1e293b"
        for c in colors:
            r, g, b = c.rgb
            if (r + g + b) / 3 < 80:  # Very dark color
                text_color = c.hex_value
                break
        
        palette = ColorPalette.from_hex_dict({
            "primary": get_color_or_default(0, "#2563eb"),
            "secondary": get_color_or_default(1, "#7c3aed"),
            "accent": get_color_or_default(2, "#f59e0b"),
            "background": background,
            "surface": "#f8fafc",
            "text_primary": text_color,
            "text_secondary": "#64748b",
        })
        
        # Determine fonts
        fonts = extraction_result.fonts
        heading_font = "Inter"
        body_font = "Inter"
        
        for f in fonts:
            if f.is_heading:
                heading_font = f.name
                break
        
        for f in fonts:
            if not f.is_heading:
                body_font = f.name
                break
        
        # Calculate typography params
        heading_sizes = [f.avg_size_pt for f in fonts if f.is_heading and f.avg_size_pt > 0]
        body_sizes = [f.avg_size_pt for f in fonts if not f.is_heading and f.avg_size_pt > 0]
        
        base_size = sum(body_sizes) / len(body_sizes) if body_sizes else 14.0
        heading_size = sum(heading_sizes) / len(heading_sizes) if heading_sizes else 28.0
        scale_ratio = (heading_size / base_size) ** 0.25 if base_size > 0 else 1.25
        
        typography = TypographyParams(
            heading_font=heading_font,
            body_font=body_font,
            heading_weight=FontWeight.BOLD,
            body_weight=FontWeight.REGULAR,
            base_size_pt=base_size,
            scale_ratio=max(1.1, min(1.5, scale_ratio)),
        )
        
        # Layout physics from metrics
        metrics = extraction_result.layout_metrics
        layout = LayoutPhysics(
            density=metrics.content_density,
            corner_radius=8.0,  # Default, can't easily extract from PPT
            margin_factor=metrics.avg_margin_ratio,
            shadow_intensity=0.1,
        )
        
        return StyleGene(
            gene_id=f"extracted_{gene_name.replace(' ', '_').lower()}",
            name=gene_name,
            typography=typography,
            palette=palette,
            layout=layout,
            visual_style_prompt=f"professional presentation, {heading_font} font, clean layout",
            tags=["extracted", "reference", "custom"],
        )


# Convenience functions
def analyze_pptx(file_path: str) -> ExtractionResult:
    """Quick function to analyze a PPTX file."""
    analyzer = ReferenceAnalyzer()
    return analyzer.analyze_presentation(file_path)


def extract_style_from_pptx(file_path: str, name: str = "extracted") -> StyleGene:
    """Extract a StyleGene from a reference PPTX file."""
    analyzer = ReferenceAnalyzer()
    result = analyzer.analyze_presentation(file_path)
    return analyzer.extract_style_gene(result, name)
