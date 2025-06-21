import os
import openai
from pptx import Presentation
from pptx.util import Inches

openai.api_key = os.getenv('OPENAI_API_KEY')


def generate_outline(title, num_slides):
    prompt = (
        f"Create an outline with {num_slides} slide titles for a PowerPoint presentation titled '{title}'."
    )
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}]
    )
    raw_lines = response.choices[0].message.content.split('\n')
    slide_titles = [line.strip().split('. ', 1)[-1] for line in raw_lines if line.strip()]
    return slide_titles


def generate_slide_content(slide_title):
    prompt = (
        f"Generate 3 to 5 bullet points for a PowerPoint slide titled '{slide_title}'."
        " Keep it concise and clear."
    )
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}]
    )
    raw_lines = response.choices[0].message.content.split('\n')
    bullet_points = [line.strip('-• ').strip() for line in raw_lines if line.strip()]
    return bullet_points


def create_presentation(title, num_slides):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated using OpenAI"

    slide_titles = generate_outline(title, num_slides)
    for slide_title in slide_titles:
        bullet_points = generate_slide_content(slide_title)
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = slide_title
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    output_path = os.path.join('presentations', f"{title.replace(' ', '_')}.pptx")
    os.makedirs('presentations', exist_ok=True)
    prs.save(output_path)
    return output_path
