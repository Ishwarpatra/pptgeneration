import os
import argparse
from pptx import Presentation
import openai


def extract_text_and_styles(ppt_path):
    """Extract text content and basic style info from a PPTX file."""
    prs = Presentation(ppt_path)
    fonts = set()
    colors = set()
    slides = []
    layouts = []

    for slide in prs.slides:
        slide_info = {"texts": []}
        layouts.append(slide.slide_layout.name)
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_info["texts"].append(shape.text)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        if run.font.color and run.font.color.type == 1:
                            colors.add(str(run.font.color.rgb))
        slides.append(slide_info)

    return {"slides": slides, "fonts": list(fonts), "colors": list(colors), "layouts": layouts}


def replicate_presentation(ppt_paths, output_path):
    """Create a new presentation with text content copied from reference PPTX files."""
    if not ppt_paths:
        raise ValueError("No input paths provided")

    first = Presentation(ppt_paths[0])
    new_prs = Presentation()
    new_prs.slide_height = first.slide_height
    new_prs.slide_width = first.slide_width

    for path in ppt_paths:
        ref = Presentation(path)
        for slide in ref.slides:
            layout = new_prs.slide_layouts[0]
            new_slide = new_prs.slides.add_slide(layout)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    textbox.text = shape.text

    new_prs.save(output_path)
    return output_path


def suggest_headline(text):
    """Return a short headline suggestion for given text using OpenAI if available."""
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return None
    openai.api_key = api_key
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": f"Suggest a short headline for: {text}"}]
        )
        return response.choices[0].message.content.strip()
    except Exception:
        return None


def main():
    parser = argparse.ArgumentParser(description="Analyze and replicate PowerPoint presentations")
    parser.add_argument("input", nargs='+', help="Path(s) to reference PPTX files")
    parser.add_argument("-o", "--output", default="replicated.pptx", help="Output PPTX file path")
    parser.add_argument("--suggest", action="store_true", help="Use AI to suggest slide headlines")
    args = parser.parse_args()

    analyses = []
    for p in args.input:
        info = extract_text_and_styles(p)
        analyses.append(info)
        print(f"Analyzed {p}: {len(info['slides'])} slides found")

    result = replicate_presentation(args.input, args.output)
    print(f"Replicated presentation saved to {result}")

    if args.suggest:
        for analysis in analyses:
            for slide in analysis['slides']:
                for text in slide['texts']:
                    suggestion = suggest_headline(text)
                    if suggestion:
                        print(f"Suggestion for '{text[:30]}...': {suggestion}")

if __name__ == "__main__":
    main()
