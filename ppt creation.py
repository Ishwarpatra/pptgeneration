import os
import openai
from pptx import Presentation
from pptx.util import Inches

# ----------------------------
# Configuration
# ----------------------------

# Set your OpenAI API key via environment variable for security
openai.api_key = os.getenv("OPENAI_API_KEY")

# ----------------------------
# Generate Slide Outline
# ----------------------------

def generate_outline(title, num_slides):
    prompt = f"Create an outline with {num_slides} slide titles for a PowerPoint presentation titled '{title}'."
    
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}]
    )
    
    raw_lines = response.choices[0].message.content.split('\n')
    slide_titles = [line.strip().split('. ', 1)[-1] for line in raw_lines if line.strip()]
    return slide_titles

# ----------------------------
# Generate Slide Content
# ----------------------------

def generate_slide_content(slide_title):
    prompt = f"Generate 3 to 5 bullet points for a PowerPoint slide titled '{slide_title}'. Keep it concise and clear."
    
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}]
    )
    
    raw_lines = response.choices[0].message.content.split('\n')
    bullet_points = [line.strip('-• ').strip() for line in raw_lines if line.strip()]
    return bullet_points

# ----------------------------
# Create PowerPoint File
# ----------------------------

def create_presentation(title, num_slides):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated using OpenAI GPT-3.5 Turbo"

    # Generate outline and content
    slide_titles = generate_outline(title, num_slides)

    for slide_title in slide_titles:
        bullet_points = generate_slide_content(slide_title)
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = slide_title
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # Save file
    output_path = "ai_generated_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

# ----------------------------
# Example Usage
# ----------------------------

if __name__ == "__main__":
    # Change the title and number of slides as needed
    create_presentation("The Impact of Artificial Intelligence on Modern Education", num_slides=6)
